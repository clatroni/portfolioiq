# -*- coding: utf-8 -*-
"""Build PortfolioIQ_Enerwave_Presentation.pptx in the HOUSE format (20 x 11.25),
matching Enerwave_Deloitte_Portfolio_IQ.pptx: Aptos, green #86BC25 accent,
brand lockup top-left, dark-green breadcrumb, 54pt title with green period,
light/dark output cards, dark section dividers with giant numerals.

Positions Portfolio IQ as a NEW AI layer on top of the already-delivered reporting solution.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
BASE = os.path.join(HERE, "output", "Enerwave_Deloitte_Portfolio_IQ.pptx")  # inherit theme/master
OUT  = os.path.join(HERE, "output", "PortfolioIQ_Enerwave_Presentation.pptx")

# ---- palette (from house deck) ----
BLACK  = RGBColor(0x22, 0x22, 0x22)
INK    = RGBColor(0x1C, 0x1C, 0x1C)
GREEN  = RGBColor(0x86, 0xBC, 0x25)
DGREEN = RGBColor(0x26, 0x89, 0x0D)   # label / breadcrumb green
DKCARD = RGBColor(0x1C, 0x3D, 0x26)   # dark green card bg
GRAY   = RGBColor(0x53, 0x56, 0x5A)
MUTE   = RGBColor(0x75, 0x78, 0x7B)
LINE   = RGBColor(0xE6, 0xE6, 0xE6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LCARD  = RGBColor(0xF7, 0xF9, 0xF5)   # light card bg
RED    = RGBColor(0xC0, 0x00, 0x00)
AMBER  = RGBColor(0xD9, 0x8A, 0x00)
PUREBLK = RGBColor(0x00, 0x00, 0x00)
F = "Aptos"

prs = Presentation(BASE)
# strip all existing slides, keep masters/layouts/theme
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(qn("r:id"))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# pick the blank "DEFAULT" layout
LAYOUT = None
for lay in prs.slide_layouts:
    if lay.name == "DEFAULT":
        LAYOUT = lay; break
if LAYOUT is None:
    LAYOUT = prs.slide_layouts[1]

SW, SH = prs.slide_width/914400.0, prs.slide_height/914400.0  # 20 x 11.25

def slide():
    s = prs.slides.add_slide(LAYOUT)
    # strip any inherited placeholders for a clean canvas
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    return s

def tb(s, l, t, w, h, anchor=None):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    if anchor is not None: tf.vertical_anchor = anchor
    return tf

def _st(p, size, bold, color, align, font=F, sa=None, sb=None):
    p.font.size = Pt(size); p.font.bold = bold; p.font.name = font
    p.font.color.rgb = color; p.alignment = align
    if sa is not None: p.space_after = Pt(sa)
    if sb is not None: p.space_before = Pt(sb)

def p0(tf, text, size=14, bold=False, color=GRAY, align=PP_ALIGN.LEFT, font=F, sa=None):
    p = tf.paragraphs[0]; p.text = str(text); _st(p, size, bold, color, align, font, sa); return p

def ap(tf, text, size=14, bold=False, color=GRAY, align=PP_ALIGN.LEFT, font=F, sa=4, sb=None):
    p = tf.add_paragraph(); p.text = str(text); _st(p, size, bold, color, align, font, sa, sb); return p

def runs_par(tf, segments, size, first=False, align=PP_ALIGN.LEFT, sa=4, sb=None):
    """segments: list of (text, bold, color)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if sa is not None: p.space_after = Pt(sa)
    if sb is not None: p.space_before = Pt(sb)
    for text, bold, color in segments:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = F; r.font.color.rgb = color
    return p

def rect(s, l, t, w, h, color, ln=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if ln is None: sp.line.fill.background()
    else: sp.line.color.rgb = ln; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False; return sp

def oval(s, l, t, w, h, color=None, ln=None, lw=1.5):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    if color is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = color
    if ln is None: sp.line.fill.background()
    else: sp.line.color.rgb = ln; sp.line.width = Pt(lw)
    sp.shadow.inherit = False; return sp

# ---------- shared chrome ----------
def brand(s):
    tf = tb(s, 0.94, 0.48, 2.0, 0.34)
    runs_par(tf, [("Portfolio IQ", True, BLACK), (".", True, GREEN)], 17, first=True)
    rect(s, 2.78, 0.52, 0.012, 0.24, LINE)
    tf = tb(s, 2.95, 0.52, 4.0, 0.30)
    p0(tf, "DELOITTE AI & DATA", size=12, bold=False, color=MUTE)

def footer(s, text="DELOITTE  ×  ENERWAVE"):
    tf = tb(s, 0.94, 10.55, 9.0, 0.32)
    p0(tf, text, size=12, bold=False, color=MUTE)

def pagenum(s, n):
    tf = tb(s, 17.9, 10.5, 1.2, 0.32)
    p0(tf, "%02d / 10" % n, size=12, bold=False, color=MUTE, align=PP_ALIGN.RIGHT)

def breadcrumb(s, text):
    tf = tb(s, 0.94, 1.46, 18.1, 0.32)
    p0(tf, text.upper(), size=15, bold=False, color=DGREEN)

def bigtitle(s, text, period=True, y=2.06, size=50, w=18.2):
    tf = tb(s, 0.94, y, w, 1.7)
    if period:
        runs_par(tf, [(text, True, BLACK), (".", True, GREEN)], size, first=True)
    else:
        p0(tf, text, size=size, bold=True, color=BLACK)
    return tf

def lead(s, segments, y=3.78, size=16, w=18.1):
    tf = tb(s, 1.04, y, w, 0.7)
    runs_par(tf, segments, size, first=True, sa=0)
    return tf

def card(s, l, t, w, h, label, title, title_sub, desc, tags, dark=False):
    bg = DKCARD if dark else LCARD
    rect(s, l, t, w, h, bg)
    rect(s, l, t, w, 0.045, GREEN)
    tcol  = WHITE if dark else BLACK
    subc  = WHITE if dark else MUTE
    descc = WHITE if dark else GRAY
    lblc  = GREEN if dark else DGREEN
    pad = 0.28
    tf = tb(s, l+pad, t+0.18, w-2*pad, 0.30)
    p0(tf, label.upper(), size=14, bold=True, color=lblc)
    tf = tb(s, l+pad, t+0.52, w-2*pad, 0.95)
    runs_par(tf, [(title, True, tcol)], 24, first=True, sa=2)
    if title_sub:
        runs_par(tf, [(title_sub, False, subc)], 15, sa=0)
    yy = t + (1.55 if title_sub else 1.25)
    if desc:
        tf = tb(s, l+pad, yy, w-2*pad, 1.0)
        p0(tf, desc, size=15, bold=False, color=descc)
        yy += 0.95
    for term, rest in tags:
        tf = tb(s, l+pad, yy, w-2*pad, 0.32)
        runs_par(tf, [("—  ", False, GREEN), (term, True, tcol), ("  "+rest, False, subc if dark else MUTE)],
                 14, first=True, sa=0)
        yy += 0.36

def divider(s, num, title, subtitle, tag):
    rect(s, 0, 0, SW, SH, PUREBLK)
    rect(s, 0, 0, SW, 0.10, GREEN)
    oval(s, 14.0, 1.5, 8.2, 8.2, ln=DKCARD, lw=2.0)
    oval(s, 16.4, 3.9, 3.4, 3.4, color=DGREEN)
    tf = tb(s, 1.0, 0.95, 12, 0.4)
    p0(tf, "DELOITTE · AI & DATA", size=18, bold=True, color=GREEN)
    tf = tb(s, 0.96, 3.0, 14, 2.4)
    p0(tf, num, size=165, bold=True, color=GREEN)
    tf = tb(s, 1.0, 5.9, 16, 1.3)
    runs_par(tf, [(title, True, WHITE), (".", True, GREEN)], 72, first=True)
    tf = tb(s, 1.0, 7.45, 13, 0.7)
    p0(tf, subtitle, size=26, bold=False, color=WHITE)
    tf = tb(s, 1.0, 10.5, 10, 0.32)
    p0(tf, tag, size=14, bold=False, color=MUTE)

def styled_table(s, headers, rows, l, t, w, col_w, rowh, hsize=16, fsize=15):
    nr, nc = len(rows)+1, len(headers)
    shp = s.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(rowh*nr))
    tbl = shp.table; tbl.first_row = True; tbl.horz_banding = False
    for i, cw in enumerate(col_w): tbl.columns[i].width = Inches(cw)
    for i, htext in enumerate(headers):
        c = tbl.cell(0, i); c.text = htext
        c.fill.solid(); c.fill.fore_color.rgb = DKCARD
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.18); c.margin_top = Inches(0.06); c.margin_bottom = Inches(0.06)
        for p in c.text_frame.paragraphs:
            _st(p, hsize, True, GREEN if i > 0 else WHITE, PP_ALIGN.LEFT)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri+1, ci); c.text = str(val)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LCARD
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.18); c.margin_top = Inches(0.05); c.margin_bottom = Inches(0.05)
            first = (ci == 0)
            for p in c.text_frame.paragraphs:
                _st(p, fsize, first, INK if first else GRAY, PP_ALIGN.LEFT)
    return tbl

def callout(s, text, t=9.95, color=GRAY, accent=GREEN):
    rect(s, 0.94, t, 0.06, 0.40, accent)
    tf = tb(s, 1.16, t-0.02, 17.6, 0.5)
    p0(tf, text, size=14, bold=False, color=color)

# =============================================================== 1 COVER (dark hero)
s = slide()
rect(s, 0, 0, SW, SH, PUREBLK)
rect(s, 0, 0, SW, 0.10, GREEN)
oval(s, 13.6, 1.2, 9.2, 9.2, ln=DKCARD, lw=2.0)
oval(s, 16.3, 3.9, 3.8, 3.8, color=DGREEN)
rect(s, 0.96, 1.02, 0.62, 0.12, GREEN)
tf = tb(s, 1.75, 0.92, 10, 0.34)
p0(tf, "DELOITTE  ×  ENERWAVE", size=15, bold=True, color=GREEN)
tf = tb(s, 0.94, 3.85, 14.5, 2.3)
runs_par(tf, [("Portfolio IQ", True, WHITE), (".", True, GREEN)], 100, first=True)
tf = tb(s, 0.96, 6.25, 14, 0.8)
p0(tf, "Your monthly portfolio review, automated.", size=30, bold=False, color=GREEN)
tf = tb(s, 0.97, 7.15, 13.6, 0.9)
p0(tf, "Written, designed and reviewed by AI. A new layer on the reporting solution already live.",
   size=18, bold=False, color=RGBColor(0xCF,0xD3,0xCC))
tf = tb(s, 0.96, 9.78, 6, 0.32)
p0(tf, "DELOITTE AI & DATA", size=13, bold=True, color=WHITE)
tf = tb(s, 0.96, 10.12, 6, 0.30)
p0(tf, "JUNE 2026", size=13, bold=False, color=RGBColor(0xCF,0xD3,0xCC))

# =============================================================== 2 POSITIONING
s = slide(); brand(s); footer(s); pagenum(s, 2)
breadcrumb(s, "Positioning")
bigtitle(s, "A new AI layer on the solution you already have")
lead(s, [("Portfolio IQ does not replace your Power BI dashboard. ", False, GRAY),
         ("The dashboard is its input.", True, BLACK)])
styled_table(s,
    ["", "Already live  (the foundation)", "What Portfolio IQ adds  (new)"],
    [
        ["Data", "Microsoft Fabric Gold layer, certified semantic model", "Same Gold layer, read-only. No new model."],
        ["Dashboard", "Power BI portfolio dashboard the team uses today", "Becomes the agent's input, it is not replaced."],
        ["Reporting", "Month-end deck assembled manually", "Executive deck written, designed and self-reviewed by AI."],
        ["Q and A", "Ad-hoc data-pull requests to analysts", "Ask the AI, answered from the same source data."],
    ],
    0.94, 4.85, 18.12, [2.6, 7.76, 7.76], rowh=1.02)
callout(s, "One Gold layer. The existing dashboard stays. Portfolio IQ sits on top and automates the reporting cycle.")

# =============================================================== 3 HOW IT WORKS
s = slide(); brand(s); footer(s); pagenum(s, 3)
breadcrumb(s, "How it works")
bigtitle(s, "Three steps, start to finish, in under a minute")
lead(s, [("An agent loop on top of the existing dashboard. No manual data slicing, no manual slide building.", False, GRAY)])
cw, ct, ch = 5.84, 4.85, 4.95
xs = [0.94, 0.94+cw+0.30, 0.94+2*(cw+0.30)]
card(s, xs[0], ct, cw, ch, "Step 01", "Read & Analyze", "From the live dashboard",
     "Starts from the existing Power BI portfolio dashboard, the same one the team uses today.",
     [("Budget health", ""), ("Schedule slippage", ""), ("Exposure", "Red / Amber"), ("Projects at risk", "")])
card(s, xs[1], ct, cw, ch, "Step 02", "Draft & Design", "Narrative and deck",
     "Writes the executive narrative, then builds the full deck in the Deloitte corporate template.",
     [("Takeaways", ""), ("Decisions needed", ""), ("Key risks", ""), ("10-slide branded deck", "")])
card(s, xs[2], ct, cw, ch, "Step 03", "Review & Approve", "The self-QA pass",
     "Re-reads the finished deck and checks every number against the source data.",
     [("Reconciled to source", ""), ("Inconsistencies flagged", ""), ("Affected slides rebuilt", ""),
      ("Only a clean deck ships", "")], dark=True)
callout(s, "Step 3, the self-QA pass, is what makes the output safe to send unedited. It is the key differentiator.")

# =============================================================== 4 DIFFERENTIATOR
s = slide(); brand(s); footer(s); pagenum(s, 4)
breadcrumb(s, "The differentiator")
bigtitle(s, "The AI checks its own work before you see it")
lead(s, [("Every figure in the deck is re-read and reconciled against the ", False, GRAY),
         ("Gold layer source", True, BLACK), (".", False, GRAY)])
card(s, 0.94, 4.85, 8.95, 5.0, "The self-QA pass", "What it does", "",
     None,
     [("Re-reads", "the finished deck end to end"),
      ("Reconciles", "every number against the source"),
      ("Flags", "inconsistencies and rebuilds the slide"),
      ("Releases", "only a deck that matches the source")], dark=True)
card(s, 10.11, 4.85, 8.95, 5.0, "Why it matters", "What you get", "",
     None,
     [("No errors", "no silent copy-paste mistakes in board reporting"),
      ("Aligned", "narrative numbers match the table numbers"),
      ("Less manual work", "removes the PMO cross-check done today"),
      ("Trustworthy", "by construction, not by spot-check")])
callout(s, "Most automation stops at draft. Portfolio IQ adds a verification step, so review time replaces rework time.")

# =============================================================== 5 THE OUTPUT
s = slide(); brand(s); footer(s); pagenum(s, 5)
breadcrumb(s, "The output")
bigtitle(s, "A ready-to-send executive deck, every cycle")
lead(s, [("Generated at month-end close in seconds, in the Deloitte corporate template, ", False, GRAY),
         ("no manual edits", True, BLACK), (".", False, GRAY)])
card(s, 0.94, 4.85, 8.95, 5.0, "Output · status & narrative", "Status & narrative", "",
     None,
     [("Portfolio Snapshot", "scope, key metrics, period-on-period"),
      ("Schedule Health", "delays, critical path risk, milestones"),
      ("Budget Health", "variance, burn-rate, forecast to complete")])
card(s, 10.11, 4.85, 8.95, 5.0, "Output · decisions & risk", "Decisions & risk", "",
     None,
     [("Leadership Decisions", "escalations and open actions surfaced"),
      ("Projects at Risk", "AI-ranked issues with narrative context"),
      ("Key Risks", "the few items leadership must act on")], dark=True)
callout(s, "The same six sections every month, built from the live data, so the cadence never slips.")

# =============================================================== 6 ILLUSTRATIVE SNAPSHOT
s = slide(); brand(s); footer(s); pagenum(s, 6)
breadcrumb(s, "Illustrative output")
bigtitle(s, "What a generated snapshot looks like")
lead(s, [("Illustrative figures from the demo. ", False, GRAY),
         ("Mock data, not actual Enerwave portfolio data.", True, RED)])
kpis = [("201","Projects in scope",BLACK), ("100","Active",BLACK), ("€39.4M","FY budget",BLACK),
        ("11","Flagged (Red / Amber)",RED), ("€5.9M","Exposure",AMBER), ("58","No baseline",RED)]
kw, kt = 2.85, 4.85
for i,(v,lab,col) in enumerate(kpis):
    l = 0.94 + i*(kw+0.196)
    rect(s, l, kt, kw, 1.75, LCARD)
    rect(s, l, kt, kw, 0.05, GREEN)
    tf = tb(s, l+0.08, kt+0.30, kw-0.16, 0.8, anchor=MSO_ANCHOR.MIDDLE)
    p0(tf, v, size=40, bold=True, color=col, align=PP_ALIGN.CENTER)
    tf = tb(s, l+0.08, kt+1.18, kw-0.16, 0.5, anchor=MSO_ANCHOR.TOP)
    p0(tf, lab, size=13, bold=False, color=GRAY, align=PP_ALIGN.CENTER)
# takeaway card (dark)
tcard_t = 6.95
rect(s, 0.94, tcard_t, 18.12, 2.55, DKCARD)
rect(s, 0.94, tcard_t, 18.12, 0.05, GREEN)
tf = tb(s, 1.20, tcard_t+0.22, 17.6, 0.4)
p0(tf, "AI-WRITTEN KEY TAKEAWAY", size=14, bold=True, color=GREEN)
tf = tb(s, 1.20, tcard_t+0.66, 17.6, 1.7)
p0(tf, "Budget discipline remains strong, 96% of active projects are within budget.",
   size=18, bold=False, color=WHITE, sa=8)
ap(tf, "Schedule baseline data is the critical gap. 81% of the active portfolio is delayed or lacks "
       "schedule visibility, with 58 projects missing baseline dates.",
   size=18, bold=False, color=WHITE, sa=0)
callout(s, "Illustrative only. On delivery these are your live Gold-layer figures for the selected reporting period.",
        color=MUTE)

# =============================================================== 7 ASK THE AI
s = slide(); brand(s); footer(s); pagenum(s, 7)
breadcrumb(s, "Ask the AI")
bigtitle(s, "Leadership can ask the portfolio anything")
lead(s, [("A conversational assistant on the ", False, GRAY),
         ("same source data", True, BLACK), (" behind the deck.", False, GRAY)])
card(s, 0.94, 4.85, 8.95, 5.0, "Indicative questions", "Ask in plain language", "",
     None,
     [("This vs last", "“Active projects this month vs last?”"),
      ("Drivers", "“Which owners have the most at-risk projects, and why?”"),
      ("Exposure", "“How much budget is at risk this period?”"),
      ("Follow-up", "“Break exposure down by category.”")])
card(s, 10.11, 4.85, 8.95, 5.0, "How it answers", "Grounded and governed", "",
     None,
     [("Writes the query", "against the certified model"),
      ("Respects RLS", "row-level security on sensitive data"),
      ("Matches the dashboard", "returns the same figures, exactly"),
      ("Keeps context", "follow-ups need no re-explaining")], dark=True)
callout(s, "Same governed source as the deck, so the answer a leader gets is the answer the report would show.")

# =============================================================== 8 ARCHITECTURE
s = slide(); brand(s); footer(s); pagenum(s, 8)
breadcrumb(s, "Architecture")
bigtitle(s, "Built on what you already have")
lead(s, [("Microsoft Fabric and Azure AI. ", False, GRAY),
         ("No new license, no new data model.", True, BLACK)])
row_t, bh = 5.4, 3.2
def flowbox(l, w, head, lines, tag, new=False):
    bg = DKCARD if new else LCARD
    rect(s, l, row_t, w, bh, bg)
    rect(s, l, row_t, w, 0.05, GREEN)
    hc = WHITE if new else BLACK
    lc = WHITE if new else GRAY
    tf = tb(s, l+0.28, row_t+0.28, w-0.56, 0.5)
    p0(tf, head, size=20, bold=True, color=hc)
    tf = tb(s, l+0.28, row_t+0.95, w-0.56, 1.6)
    for i, ln in enumerate(lines):
        if i == 0: p0(tf, ln, size=14, bold=False, color=lc, sa=4)
        else: ap(tf, ln, size=14, bold=False, color=lc, sa=4)
    badge = tb(s, l+0.28, row_t+bh-0.5, w-0.56, 0.32)
    p0(badge, "NEW" if new else "EXISTING", size=12, bold=True, color=GREEN if new else MUTE)
w1, w2, w3, w4, g = 3.9, 4.5, 4.4, 4.1, 0.40
l1 = 0.94; l2 = l1+w1+g; l3 = l2+w2+g; l4 = l3+w3+g
flowbox(l1, w1, "Sources", ["SAP / PCS", "Excel and files"], "EXISTING")
flowbox(l2, w2, "Microsoft Fabric", ["Bronze · Silver · Gold", "Certified semantic model", "Power BI dashboard"], "EXISTING")
flowbox(l3, w3, "Portfolio IQ agent", ["Azure AI", "Read · analyze", "draft · self-review"], "NEW", new=True)
flowbox(l4, w4, "Outputs", ["Executive PPTX", "Ask the AI assistant"], "NEW", new=True)
for (la, lb) in [(l1+w1, l2), (l2+w2, l3), (l3+w3, l4)]:
    ar = tb(s, la-0.05, row_t+bh/2-0.35, lb-la+0.10, 0.7, anchor=MSO_ANCHOR.MIDDLE)
    p0(ar, "→", size=30, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
callout(s, "One source of truth, governed end to end, running on your existing Fabric capacity.")

# =============================================================== 9 WHY IT MATTERS
s = slide(); brand(s); footer(s); pagenum(s, 9)
breadcrumb(s, "Why it matters")
bigtitle(s, "From days of assembly to seconds of review")
lead(s, [("The reporting cycle stays. The manual effort and the cross-check risk do not.", False, GRAY)])
styled_table(s,
    ["", "Today", "With Portfolio IQ"],
    [
        ["Month-end deck", "Assembled by hand, hours of effort", "Auto-generated in seconds, ready to send"],
        ["Accuracy", "Numbers cross-checked manually", "Self-QA reconciles every figure to source"],
        ["Q and A", "Ad-hoc requests queue with analysts", "Self-service answers in seconds"],
        ["Consistency", "Deck and dashboard can drift apart", "One governed source, always aligned"],
        ["PMO focus", "Time spent producing the report", "Time spent on the decisions in it"],
    ],
    0.94, 4.85, 18.12, [3.2, 7.46, 7.46], rowh=0.92)
callout(s, "Same governance, same template, same data. Less production, more decision-making.")

# =============================================================== 10 CLOSING
s = slide()
divider(s, "", "Portfolio IQ", "Your monthly portfolio review, automated.",
        "DELOITTE  ×  ENERWAVE")
# replace giant numeral area with a closing line
tf = tb(s, 1.0, 3.2, 16, 0.6)
p0(tf, "THE NEXT STEP", size=18, bold=True, color=GREEN)
tf = tb(s, 1.0, 8.5, 16, 0.6)
p0(tf, "Built on Microsoft Fabric and Azure AI.  ·  A new layer on the solution already delivered.",
   size=16, bold=False, color=GREEN)

prs.save(OUT)
print("SAVED:", OUT)
print("SLIDES:", len(prs.slides), "| DIMS %.2f x %.2f" % (SW, SH))
