"""patch_period_pptxs.py — patch Dec/Feb PPTXs with per-period KPI numbers.

Run AFTER build_period_pptxs.py (which swaps period labels). This script
opens each variant with python-pptx and replaces the headline KPI values
on slide 2 (snapshot) so the numbers match each period's actual data.

What it patches (slide 2):
  - 5 KPI tile big numbers + their sub-texts
  - 3 budget-status insight lines
  - Subtitle "High-level status across N projects as of <period>"
  - Schedule + budget card details

What it does NOT patch (deliberate scope cut):
  - Slide 3 bar chart (chart XML + embedded xlsx — complex)
  - Slide 4 donut (similar)
  - Slide 6 top-5 at-risk table (would need to re-rank projects per period)
  - Slide 8 full red/amber list
  - Slide 9 timeline
  - LLM-narrated slides 5, 7 (text stays as Jan placeholder)
  - Most narrative paragraphs (left as Jan to keep the demo simple)

The result: cover + headline KPI tiles + budget insights match each period.
Tables, charts, and narrative still reflect Jan baseline.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.chart.data import CategoryChartData

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

BASE   = Path(__file__).resolve().parent.parent
PUBLIC = BASE / "website" / "public"
DATA   = PUBLIC / "data"

# Map period key -> JSON path so we can load period-specific chart/table data
PERIOD_JSON = {
    "Dec": DATA / "2025-12.json",
    "Jan": DATA / "2026-01.json",
    "Feb": DATA / "2026-02.json",
}

# Per-period values, computed from the pipeline JSONs
PERIOD_VALUES = {
    "Jan": {
        # Jan is the colleague's original render — patches here CORRECT the
        # rag count (was 13 incl. completed projects; truth is 11 active-flagged).
        # All other Jan KPIs are already correct in his original.
        "exposure_line":      "€5.9M exposure across 11 Red/Amber flagged projects.",
        # Takeaway: the colleague's LLM-generated text references 13; correct to 11.
        "takeaway_verdict":   "Portfolio performance shows mixed results with key risks identified.",
        "takeaway_detail":    "81% of active portfolio shows schedule risk: 22 major delays and 12 overdue projects. Budget adherence remains strong at 96%, yet €5.9M exposure across 11 Red/Amber projects requires attention.",
        "takeaway_action":    "Immediate action required: prioritize overdue and delayed projects to mitigate risks and ensure schedule recovery.",
    },
    "Dec": {
        "completed":          "24",
        "overdue":            "4",
        "ytd_pct":            "3.3%",
        "active_pending_sub": "96 active · 105 pending",
        "completed_sub":      "of 96 active",
        "overdue_sub":        "4% of active",
        "budget_sub":         "€18.0M active · €21.4M pending",
        "ytd_sub":            "€1.3M of €39.4M (PR+PO)",
        "subtitle":           "High-level status across 201 projects as of December 2025",
        "within_budget_line": "98% within budget — strong fiscal discipline.",
        "overrun_line":       "2 projects with Budget Overrun Risk (up to 280% of FY2025-2026 Budget).",
        "exposure_line":      "€0.9M exposure across 10 Red/Amber flagged projects.",
        "schedule_headline":  "83% of active portfolio is delayed or lacking schedule visibility.",
        "slide3_subtitle":    "Schedule distribution across 96 active projects and key risk areas",
        # Takeaway: portfolio entered the FY with a healthy baseline
        "takeaway_verdict":   "Portfolio entered FY2026 with a healthy baseline.",
        "takeaway_detail":    "10 projects flagged Red/Amber for €0.9M total exposure. Schedule adherence at 17% of active portfolio (83% delayed or no baseline) is the main concern. Budget discipline strong: 98% within budget.",
        "takeaway_action":    "Focus on closing the schedule-baseline gap; 58 projects still lack a baseline.",
    },
    "Feb": {
        "completed":          "26",
        "overdue":            "12",
        "ytd_pct":            "25.7%",
        "active_pending_sub": "98 active · 101 pending",
        "completed_sub":      "of 98 active",
        "overdue_sub":        "12% of active",
        "budget_sub":         "€21.5M active · €18.0M pending",
        "ytd_sub":            "€10.1M of €39.4M (PR+PO)",
        "subtitle":           "High-level status across 201 projects as of February 2026",
        "within_budget_line": "97% within budget — strong fiscal discipline.",
        "overrun_line":       "5 projects with Budget Overrun Risk (up to 300% of FY2026 Budget).",
        "exposure_line":      "€9.2M exposure across 12 Red/Amber flagged projects.",
        "schedule_headline":  "81% of active portfolio is delayed or lacking schedule visibility.",
        "slide3_subtitle":    "Schedule distribution across 98 active projects and key risk areas",
        # Takeaway: portfolio risk profile escalating, exposure tripled in 2 months
        "takeaway_verdict":   "Portfolio risk profile escalating; exposure 10× since December.",
        "takeaway_detail":    "12 flagged projects now carry €9.2M exposure (vs €0.9M in Dec). Steam Turbine Generator rewinding emerged as a new €3M risk. Schedule slippage stable at 81% but exposure is growing through bigger-ticket flags.",
        "takeaway_action":    "Immediate action required: validate Steam Turbine Generator scope and de-risk Q1 schedule recovery before further escalation.",
    },
}

# Strings that uniquely identify the Jan content we're replacing
JAN_MARKERS = {
    "completed_big":      "25",
    "overdue_big":        "12",
    "ytd_pct_big":        "16.5%",
    "active_pending_sub": "100 active · 101 pending",
    "completed_sub":      "of 100 active",
    "overdue_sub":        "12% of active",
    "budget_sub":         "€21.5M active · €18.0M pending",
    "ytd_sub":            "€6.5M of €39.4M (PR+PO)",
    "subtitle_prefix":    "High-level status across 201 projects as of",
    "within_budget_line": "96% within budget",
    "overrun_line":       "4 projects with Budget Overrun Risk",
    "exposure_line":      "€5.9M exposure across 13 Red/Amber",   # also Jan: rewrite to 11
    "schedule_headline":  "81% of active portfolio is delayed",
    "slide3_subtitle":    "Schedule distribution across 100 active projects",
    # Takeaway paragraphs (LLM-narrated in Jan deck — re-template per period)
    "takeaway_verdict":   "Portfolio performance shows mixed",     # unique prefix
    "takeaway_detail":    "of projects are aligned with schedules",  # unique fragment
    "takeaway_action":    "Immediate action required: Prioritize",
}


def set_text_preserving_format(tf, new_text: str) -> bool:
    """Replace the text in a text_frame, preserving the first run's font/color/size."""
    if not tf.paragraphs: return False
    para = tf.paragraphs[0]
    if not para.runs: return False
    # Wipe runs 2..N (keep first run for formatting)
    for run in para.runs[1:]:
        run.text = ""
    para.runs[0].text = new_text
    # Wipe additional paragraphs if any
    for extra_para in tf.paragraphs[1:]:
        for run in extra_para.runs:
            run.text = ""
    return True


def replace_substring(tf, old: str, new: str) -> bool:
    """Find a run containing the old substring, do an in-run replace."""
    for para in tf.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                return True
    return False


def set_cell_text(cell, new_text: str) -> None:
    """Set a table cell's text while preserving first-run formatting (font/colour)."""
    tf = cell.text_frame
    if not tf.paragraphs:
        cell.text = new_text
        return
    p = tf.paragraphs[0]
    if not p.runs:
        p.text = new_text
    else:
        p.runs[0].text = new_text
        for r in p.runs[1:]: r.text = ""
    for extra in tf.paragraphs[1:]:
        for r in extra.runs: r.text = ""


def patch_chart_data(chart, categories, values, series_name="Series 1") -> None:
    """Replace chart's underlying data. Works for bar + doughnut + similar."""
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series(series_name, values)
    chart.replace_data(cd)


def _highlights_for(p) -> str:
    """Deterministic 'Highlights / Action' string per project row.
    Mirrors the colleague's LLM logic for the common patterns without calling the LLM."""
    pct  = p.get("pct", 0)
    bud  = p.get("budget", "")
    goal = p.get("go_live", "")
    # Data-quality pattern: 100% complete with €0 budget — RAG never updated
    if pct == 100 and bud in ("€0", "€0k", "€0M"):
        return "100% complete with stale RAG · Action: Update RAG to Green"
    if pct == 0:
        return f"0% complete · {bud} budget · Action: Confirm scope and start"
    if pct >= 80:
        return f"{pct}% complete · {bud} budget · Action: Push to closeout"
    return f"{pct}% complete · {bud} · target {goal} · Action: Escalate to PM"


def patch_top5_table(table, projects) -> int:
    """Replace rows 1-5 of slide 6 top-5 table with this period's top-5 projects.
    Table cols: RAG | Project | Owner | Target Go-Live | CY Budget | Progress | Highlights & Action.
    Returns number of data rows actually filled."""
    rows = list(table.rows)
    DATA_ROWS = 5  # rows 1..5 (row 0 is header)
    filled = 0
    for i in range(DATA_ROWS):
        if i + 1 >= len(rows): break
        cells = list(rows[i + 1].cells)
        if i < len(projects):
            p = projects[i]
            set_cell_text(cells[0], p["rag"])
            set_cell_text(cells[1], p["project"])
            set_cell_text(cells[2], p["owner"])
            set_cell_text(cells[3], p["go_live"])
            set_cell_text(cells[4], p["budget"])
            set_cell_text(cells[5], f"{p['pct']}%")
            set_cell_text(cells[6], _highlights_for(p))
            filled += 1
        else:
            for c in cells: set_cell_text(c, "")
    return filled


def patch_appendix_table(table, projects) -> int:
    """Replace rows 1..N of slide 8 appendix table with all this period's projects.
    Table cols: RAG | Code | Project | Owner | Go-Live | Budget | % | Highlights / Action.
    Blanks any rows past len(projects). Returns rows filled."""
    rows = list(table.rows)
    capacity = len(rows) - 1  # data row capacity (excluding header)
    filled = 0
    for i in range(capacity):
        cells = list(rows[i + 1].cells)
        if i < len(projects):
            p = projects[i]
            set_cell_text(cells[0], p["rag"])
            set_cell_text(cells[1], p["code"])
            set_cell_text(cells[2], p["project"])
            set_cell_text(cells[3], p["owner"])
            set_cell_text(cells[4], p["go_live"])
            set_cell_text(cells[5], p["budget"])
            set_cell_text(cells[6], f"{p['pct']}%")
            set_cell_text(cells[7], _highlights_for(p))
            filled += 1
        else:
            for c in cells: set_cell_text(c, "")
    return filled


def replace_paragraph_text(tf, marker: str, new_text: str) -> bool:
    """Find the PARAGRAPH within tf whose joined text contains marker, then
    replace that paragraph's content (preserving first run's formatting).
    Use this for multi-paragraph text frames (e.g. the 3-paragraph takeaway)
    where set_text_preserving_format would obliterate the other paragraphs."""
    for para in tf.paragraphs:
        para_text = "".join(run.text for run in para.runs)
        if marker in para_text:
            for run in para.runs[1:]:
                run.text = ""
            if para.runs:
                para.runs[0].text = new_text
            return True
    return False


def patch_pptx(pptx_path: Path, period_key: str):
    """Open, patch slide 2 + slide 3 subtitle, save in place.
    Jan deck only gets exposure_line + takeaway corrections (count was 13, truth = 11).
    Dec/Feb decks get the full KPI swap.
    """
    vals = PERIOD_VALUES[period_key]
    prs  = Presentation(str(pptx_path))
    replaced = []
    is_jan = (period_key == "Jan")

    # ─── Slide 2 (index 1) — Portfolio Snapshot ─────────────────────
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if not shape.has_text_frame: continue
        t = shape.text_frame.text.strip()

        # For non-Jan periods, do the full KPI swap:
        if not is_jan:
            if   t == JAN_MARKERS["completed_big"]:
                if set_text_preserving_format(shape.text_frame, vals["completed"]):
                    replaced.append(f"KPI: 25 -> {vals['completed']}")
                continue
            elif t == JAN_MARKERS["overdue_big"]:
                if set_text_preserving_format(shape.text_frame, vals["overdue"]):
                    replaced.append(f"KPI: 12 -> {vals['overdue']}")
                continue
            elif t == JAN_MARKERS["ytd_pct_big"]:
                if set_text_preserving_format(shape.text_frame, vals["ytd_pct"]):
                    replaced.append(f"KPI: 16.5% -> {vals['ytd_pct']}")
                continue
            elif t == JAN_MARKERS["active_pending_sub"]:
                if set_text_preserving_format(shape.text_frame, vals["active_pending_sub"]):
                    replaced.append(f"sub: active/pending")
                continue
            elif t == JAN_MARKERS["completed_sub"]:
                if set_text_preserving_format(shape.text_frame, vals["completed_sub"]):
                    replaced.append(f"sub: 'of N active'")
                continue
            elif t == JAN_MARKERS["overdue_sub"]:
                if set_text_preserving_format(shape.text_frame, vals["overdue_sub"]):
                    replaced.append(f"sub: overdue pct")
                continue
            elif t == JAN_MARKERS["budget_sub"]:
                if set_text_preserving_format(shape.text_frame, vals["budget_sub"]):
                    replaced.append(f"sub: budget split")
                continue
            elif t == JAN_MARKERS["ytd_sub"]:
                if set_text_preserving_format(shape.text_frame, vals["ytd_sub"]):
                    replaced.append(f"sub: YTD spend")
                continue
            elif JAN_MARKERS["within_budget_line"] in t:
                if set_text_preserving_format(shape.text_frame, vals["within_budget_line"]):
                    replaced.append(f"budget insight 1 (within)")
                continue
            elif JAN_MARKERS["overrun_line"] in t:
                if set_text_preserving_format(shape.text_frame, vals["overrun_line"]):
                    replaced.append(f"budget insight 2 (overrun)")
                continue
            elif JAN_MARKERS["schedule_headline"] in t:
                if set_text_preserving_format(shape.text_frame, vals["schedule_headline"]):
                    replaced.append(f"schedule headline")
                continue

        # Patches applied to ALL periods (Jan + Dec + Feb):
        if JAN_MARKERS["exposure_line"] in t:
            if set_text_preserving_format(shape.text_frame, vals["exposure_line"]):
                replaced.append(f"exposure line ({vals['exposure_line'][:30]}...)")
        elif any(m in t for m in (JAN_MARKERS["takeaway_verdict"],
                                  JAN_MARKERS["takeaway_detail"],
                                  JAN_MARKERS["takeaway_action"])):
            # Takeaway is 3 paragraphs in ONE text frame — replace paragraph-by-paragraph
            if replace_paragraph_text(shape.text_frame, JAN_MARKERS["takeaway_verdict"], vals["takeaway_verdict"]):
                replaced.append(f"takeaway verdict")
            if replace_paragraph_text(shape.text_frame, JAN_MARKERS["takeaway_detail"], vals["takeaway_detail"]):
                replaced.append(f"takeaway detail")
            if replace_paragraph_text(shape.text_frame, JAN_MARKERS["takeaway_action"], vals["takeaway_action"]):
                replaced.append(f"takeaway action")

    # ─── Slide 3 (index 2) — Schedule Health: subtitle + bar chart ──
    slide3 = prs.slides[2]
    if not is_jan:
        for shape in slide3.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if JAN_MARKERS["slide3_subtitle"] in t:
                    if set_text_preserving_format(shape.text_frame, vals["slide3_subtitle"]):
                        replaced.append("slide3 subtitle")
    # Load period JSON for chart + table patching (all periods, incl. Jan for table-row-count correction)
    period_data = json.load(open(PERIOD_JSON[period_key], encoding="utf-8"))
    # Patch the bar chart on slide 3
    for shape in slide3.shapes:
        if shape.has_chart:
            cats = [s["label"] for s in period_data["schedule"]]
            vals_ = [s["value"]  for s in period_data["schedule"]]
            try:
                patch_chart_data(shape.chart, cats, vals_, series_name="Schedule Health")
                replaced.append(f"slide3 bar chart: {vals_}")
            except Exception as e:
                replaced.append(f"slide3 chart FAILED: {e}")
            break

    # ─── Slide 4 (index 3) — Budget Health: donut chart ─────────────
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_chart:
            cats = [b["label"]   for b in period_data["budget"]]
            vals_ = [b["value_m"] for b in period_data["budget"]]
            try:
                patch_chart_data(shape.chart, cats, vals_, series_name="Budget Category")
                replaced.append(f"slide4 donut: {vals_}")
            except Exception as e:
                replaced.append(f"slide4 donut FAILED: {e}")
            break

    # ─── Slide 6 (index 5) — Top 5 at-risk projects table ───────────
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_table:
            n = patch_top5_table(shape.table, period_data["projects"][:5])
            replaced.append(f"slide6 top-5 table: {n} rows filled")
            break

    # ─── Slide 8 (index 7) — Full red/amber appendix table ──────────
    slide8 = prs.slides[7]
    for shape in slide8.shapes:
        if shape.has_table:
            n = patch_appendix_table(shape.table, period_data["projects"])
            replaced.append(f"slide8 appendix: {n} rows filled (capacity 13)")
            break

    # ─── Slide 3 risk bullets (3) — each bullet is its OWN shape starting with "●" ─
    bullets = period_data.get("schedule_bullets", [])
    if bullets and len(bullets) == 3:
        bullet_shapes = []
        for shape in slide3.shapes:
            if not shape.has_text_frame: continue
            t = shape.text_frame.text.strip()
            # Bullet shapes start with the bullet glyph
            if t.startswith("●") and len(shape.text_frame.paragraphs) == 1:
                bullet_shapes.append(shape)
        if len(bullet_shapes) == 3:
            for i, shape in enumerate(bullet_shapes):
                # Preserve the "●  " prefix from the original
                existing = shape.text_frame.text
                import re
                m = re.match(r"^(●\s+)", existing)
                prefix = m.group(1) if m else "●  "
                set_text_preserving_format(shape.text_frame, prefix + bullets[i])
            replaced.append(f"slide3 risk bullets ({len(bullet_shapes)} individual shapes)")
        else:
            replaced.append(f"slide3 bullets: SKIPPED (found {len(bullet_shapes)} bullet shapes, expected 3)")

    # ─── Slide 4 budget insights (4) — each insight is its OWN bullet shape ──────
    insights = period_data.get("budget_insights", [])
    if insights and len(insights) == 4:
        bullet_shapes = []
        for shape in slide4.shapes:
            if not shape.has_text_frame: continue
            t = shape.text_frame.text.strip()
            if t.startswith("●") and len(shape.text_frame.paragraphs) == 1:
                bullet_shapes.append(shape)
        if len(bullet_shapes) == 4:
            for i, shape in enumerate(bullet_shapes):
                existing = shape.text_frame.text
                import re
                m = re.match(r"^(●\s+)", existing)
                prefix = m.group(1) if m else "●  "
                set_text_preserving_format(shape.text_frame, prefix + insights[i]["bold"] + insights[i].get("tail", ""))
            replaced.append(f"slide4 budget insights ({len(bullet_shapes)} individual shapes)")
        else:
            replaced.append(f"slide4 insights: SKIPPED (found {len(bullet_shapes)} bullet shapes, expected 4)")

    # ─── Slide 5 decision cards (4) — DEFERRED ──────────────────────
    # The template uses position-based binding (each card occupies a fixed shape
    # group). Without OOXML inspection of the template (which requires the missing
    # parts/ folder from gtsimogiannis), we can't safely patch per-card text.
    # Slide 5 keeps Jan placeholder content. JSON has decision_cards array ready
    # to consume when proper rendering becomes available.
    cards = period_data.get("decision_cards", [])
    if cards:
        replaced.append(f"slide5 decision cards: JSON ready ({len(cards)} cards) but PPT patch DEFERRED — needs gtsimogiannis parts/ folder for safe per-card binding")

    # ─── Slide 7 key risks (3) ──────────────────────────────────────
    risks = period_data.get("key_risks", [])
    if risks:
        slide7 = prs.slides[6]
        for shape in slide7.shapes:
            if shape.has_table:
                tbl = shape.table
                rows = list(tbl.rows)
                # rows[0] = header. Each subsequent row = one risk.
                # Cols (per his spec): # | Type | Description / Impact | Severity | Status | Mitigation | Owner | Target
                filled = 0
                for i, risk in enumerate(risks):
                    if i + 1 >= len(rows): break
                    cells = list(rows[i + 1].cells)
                    if len(cells) >= 8:
                        # Skip col 0 (row number stays). Patch 1..7
                        set_cell_text(cells[1], risk.get("type", ""))
                        set_cell_text(cells[2], risk.get("description", ""))
                        set_cell_text(cells[3], risk.get("severity", ""))
                        set_cell_text(cells[4], risk.get("status", ""))
                        set_cell_text(cells[5], risk.get("mitigation", ""))
                        set_cell_text(cells[6], risk.get("owner", ""))
                        set_cell_text(cells[7], risk.get("target", ""))
                        filled += 1
                # Blank any extra rows
                for i in range(len(risks), len(rows) - 1):
                    cells = list(rows[i + 1].cells)
                    for c in cells: set_cell_text(c, "")
                replaced.append(f"slide7 key risks: {filled} rows filled")
                break

    prs.save(str(pptx_path))
    return replaced


def main():
    targets = [
        (PUBLIC / "PortfolioIQ_December2025.pptx", "Dec"),
        (PUBLIC / "PortfolioIQ_January2026.pptx",  "Jan"),
        (PUBLIC / "PortfolioIQ_February2026.pptx", "Feb"),
    ]
    for path, period_key in targets:
        if not path.exists():
            print(f"  SKIP {path.name} (not found — run build_period_pptxs.py first)")
            continue
        print(f"Patching {path.name} ({period_key}):")
        changes = patch_pptx(path, period_key)
        for c in changes:
            print(f"  - {c}")
        print(f"  Total: {len(changes)} patches")
        print()
    print("Done.")


if __name__ == "__main__":
    main()
