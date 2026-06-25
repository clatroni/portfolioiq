# -*- coding: utf-8 -*-
"""render_live.py — self-contained executive-deck renderer for live "Generate".

Reads a period's dashboard JSON (website/public/data/<ymd>.json) and produces the
8-slide Portfolio IQ executive deck with python-pptx, in the house format
(20 x 11.25, Deloitte palette, Aptos). Because it renders from the same JSON the
dashboard shows, the deck and the dashboard always agree.

This replaces the original jinja-parts renderer for the live endpoint: that renderer's
slide templates + source.pptx live under the git-ignored input/ tree and are not present
on this machine, so it cannot run here. This renderer has no external-asset dependency.

    python scripts/render_live.py 2026-01      # builds website/public/PortfolioIQ_January2026.pptx
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).resolve().parent.parent
DATA = BASE / "website" / "public" / "data"
PUBLIC = BASE / "website" / "public"

# ---- palette ----
GREEN  = RGBColor(0x86, 0xBC, 0x25)
DGREEN = RGBColor(0x26, 0x89, 0x0D)
DKCARD = RGBColor(0x1C, 0x3D, 0x26)
BLACK  = RGBColor(0x22, 0x22, 0x22)
INK    = RGBColor(0x1C, 0x1C, 0x1C)
GRAY   = RGBColor(0x53, 0x56, 0x5A)
MUTE   = RGBColor(0x75, 0x78, 0x7B)
LINE   = RGBColor(0xE6, 0xE6, 0xE6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LCARD  = RGBColor(0xF7, 0xF9, 0xF5)
PUREBLK = RGBColor(0x00, 0x00, 0x00)
RED    = RGBColor(0xC0, 0x00, 0x00)
AMBER  = RGBColor(0xED, 0x8B, 0x00)
OKGRN  = RGBColor(0x00, 0xB0, 0x50)
GREY   = RGBColor(0x99, 0x99, 0x99)
F = "Aptos"

RAG = {"Red": RED, "Amber": AMBER, "Green": OKGRN}
SEV = {"CRITICAL": RED, "HIGH": AMBER, "MEDIUM": RGBColor(0x00, 0x76, 0xA8), "LOW": MUTE,
       "MED": AMBER}
ACCENT = {"green": OKGRN, "amber": AMBER, "red": RED, "grey": GREY, "gray": GREY}


def _hex(s, default=GREEN):
    try:
        s = s.lstrip("#")
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except Exception:
        return default


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Emu(int(20 * 914400))
        self.prs.slide_height = Emu(int(11.25 * 914400))
        self.blank = self.prs.slide_layouts[6]
        self.SW, self.SH = 20.0, 11.25

    def slide(self):
        return self.prs.slides.add_slide(self.blank)

    def tb(self, s, l, t, w, h, anchor=None):
        box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        if anchor is not None:
            tf.vertical_anchor = anchor
        return tf

    @staticmethod
    def _st(p, size, bold, color, align, sa=None):
        p.font.size = Pt(size); p.font.bold = bold; p.font.name = F
        p.font.color.rgb = color; p.alignment = align
        if sa is not None:
            p.space_after = Pt(sa)

    def p0(self, tf, text, size=14, bold=False, color=GRAY, align=PP_ALIGN.LEFT, sa=None):
        p = tf.paragraphs[0]; p.text = str(text); self._st(p, size, bold, color, align, sa); return p

    def ap(self, tf, text, size=14, bold=False, color=GRAY, align=PP_ALIGN.LEFT, sa=4):
        p = tf.add_paragraph(); p.text = str(text); self._st(p, size, bold, color, align, sa); return p

    def runs(self, tf, segs, size, first=False, align=PP_ALIGN.LEFT, sa=4):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = align
        if sa is not None:
            p.space_after = Pt(sa)
        for text, bold, color in segs:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold; r.font.name = F; r.font.color.rgb = color
        return p

    def rect(self, s, l, t, w, h, color):
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
        sp.shadow.inherit = False; return sp

    # ---- shared chrome ----
    def brand(self, s):
        tf = self.tb(s, 0.94, 0.48, 3.0, 0.34)
        self.runs(tf, [("Portfolio IQ", True, BLACK), (".", True, GREEN)], 17, first=True)
        self.rect(s, 2.85, 0.52, 0.012, 0.24, LINE)
        tf = self.tb(s, 3.02, 0.52, 5.0, 0.3)
        self.p0(tf, "DELOITTE AI & DATA", size=12, color=MUTE)

    def footer(self, s, n):
        tf = self.tb(s, 0.94, 10.6, 9, 0.3)
        self.p0(tf, "DELOITTE  ×  ENERWAVE", size=12, color=MUTE)
        tf = self.tb(s, 17.9, 10.55, 1.2, 0.3)
        self.p0(tf, "%02d / 08" % n, size=12, color=MUTE, align=PP_ALIGN.RIGHT)

    def breadcrumb(self, s, text):
        tf = self.tb(s, 0.94, 1.4, 18.1, 0.32)
        self.p0(tf, text.upper(), size=15, color=DGREEN)

    def title(self, s, text, period=True, size=44):
        tf = self.tb(s, 0.94, 1.95, 18.2, 1.2)
        if period:
            self.runs(tf, [(text, True, BLACK), (".", True, GREEN)], size, first=True)
        else:
            self.p0(tf, text, size=size, bold=True, color=BLACK)

    def subtitle(self, s, text, y=3.2):
        tf = self.tb(s, 0.96, y, 18, 0.5)
        self.p0(tf, text, size=15, color=GRAY)

    def card_head(self, s, l, t, w, h, fill, accent=GREEN):
        self.rect(s, l, t, w, h, fill)
        self.rect(s, l, t, w, 0.05, accent)


def render(data: dict) -> Presentation:
    d = Deck()
    disp = data.get("period_display", "")
    summ = data.get("summary", {})

    # ===================================================== 1 COVER
    s = d.slide()
    d.rect(s, 0, 0, d.SW, d.SH, PUREBLK)
    d.rect(s, 0, 0, d.SW, 0.10, GREEN)
    tf = d.tb(s, 0.95, 1.0, 11, 0.34)
    d.p0(tf, "DELOITTE  ×  ENERWAVE", size=15, bold=True, color=GREEN)
    tf = d.tb(s, 0.92, 3.6, 17, 1.4)
    d.runs(tf, [("Project Portfolio", True, WHITE), (".", True, GREEN)], 60, first=True)
    tf = d.tb(s, 0.95, 5.0, 16, 0.9)
    d.p0(tf, "Executive Dashboard", size=30, color=GREEN)
    tf = d.tb(s, 0.96, 6.2, 14, 0.6)
    d.p0(tf, f"Reporting period · {disp}", size=18, color=WHITE)
    tf = d.tb(s, 0.96, 9.7, 8, 0.34)
    d.p0(tf, "AI-GENERATED EXECUTIVE REPORTING", size=12, bold=True, color=WHITE)
    tf = d.tb(s, 0.96, 10.1, 8, 0.3)
    d.p0(tf, "Written, designed and self-reviewed by AI", size=12, color=GREY)

    # ===================================================== 2 PORTFOLIO SNAPSHOT
    s = d.slide(); d.brand(s); d.footer(s, 2)
    d.breadcrumb(s, "01 · Portfolio Snapshot")
    d.title(s, "Portfolio Snapshot")
    d.subtitle(s, data.get("headline", ""))
    tiles = data.get("tiles", [])
    tw = 3.62
    for idx, tile in enumerate(tiles[:5]):
        l = 0.94 + idx * (tw + 0.13)
        d.card_head(s, l, 3.9, tw, 1.85, LCARD)
        col = ACCENT.get(tile.get("accent"), BLACK)
        tf = d.tb(s, l + 0.1, 4.12, tw - 0.2, 0.7, anchor=MSO_ANCHOR.MIDDLE)
        d.p0(tf, tile.get("value"), size=33, bold=True, color=col, align=PP_ALIGN.CENTER)
        tf = d.tb(s, l + 0.1, 4.95, tw - 0.2, 0.35)
        d.p0(tf, tile.get("label", ""), size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
        tf = d.tb(s, l + 0.1, 5.28, tw - 0.2, 0.4)
        d.p0(tf, tile.get("sub", ""), size=10.5, color=MUTE, align=PP_ALIGN.CENTER)
    # schedule + budget summary cards
    d.card_head(s, 0.94, 6.1, 9.0, 3.0, LCARD, accent=AMBER)
    tf = d.tb(s, 1.16, 6.28, 8.6, 0.4)
    d.p0(tf, "SCHEDULE HEALTH", size=14, bold=True, color=DGREEN)
    sched = data.get("schedule", [])
    tf = d.tb(s, 1.16, 6.75, 8.6, 2.2)
    for i, b in enumerate(sched):
        d.runs(tf, [(f"{b['value']:>4}  ", True, ACCENT.get(b.get('color'), INK)),
                    (b["label"], False, GRAY)], 15, first=(i == 0), sa=8)
    d.card_head(s, 10.12, 6.1, 8.94, 3.0, LCARD, accent=GREEN)
    tf = d.tb(s, 10.34, 6.28, 8.5, 0.4)
    d.p0(tf, "KEY TAKEAWAY", size=14, bold=True, color=DGREEN)
    tf = d.tb(s, 10.34, 6.75, 8.5, 2.2)
    for i, line in enumerate(data.get("takeaway", [])):
        d.ap(tf, line, size=13, color=GRAY, sa=6) if i else d.p0(tf, line, size=13, color=GRAY, sa=6)

    # ===================================================== 3 SCHEDULE HEALTH
    s = d.slide(); d.brand(s); d.footer(s, 3)
    d.breadcrumb(s, "02 · Deep Dive: Schedule Health")
    d.title(s, "Schedule Health")
    d.subtitle(s, f"Schedule distribution across {summ.get('active','-')} active projects")
    total = sum(b.get("value", 0) for b in sched) or 1
    bx, by, bw = 0.94, 4.2, 11.0
    for i, b in enumerate(sched):
        y = by + i * 1.15
        d.p0(d.tb(s, bx, y, 3.0, 0.4), b["label"], size=14, bold=True, color=INK)
        barw = max(0.15, bw * b.get("value", 0) / total)
        d.rect(s, bx + 3.1, y + 0.02, bw, 0.5, LCARD)
        d.rect(s, bx + 3.1, y + 0.02, barw, 0.5, ACCENT.get(b.get("color"), GREY))
        d.p0(d.tb(s, bx + 3.1 + bw + 0.2, y, 1.5, 0.5, anchor=MSO_ANCHOR.MIDDLE),
             str(b.get("value", 0)), size=16, bold=True, color=INK)
    d.card_head(s, 15.7, 4.2, 3.36, 4.7, DKCARD, accent=AMBER)
    tf = d.tb(s, 15.92, 4.4, 2.95, 0.4)
    d.p0(tf, "KEY SCHEDULE RISKS", size=13, bold=True, color=GREEN)
    tf = d.tb(s, 15.92, 4.9, 2.95, 3.8)
    for i, line in enumerate(data.get("schedule_bullets", [])):
        d.ap(tf, "— " + line, size=12.5, color=WHITE, sa=10) if i else d.p0(tf, "— " + line, size=12.5, color=WHITE, sa=10)

    # ===================================================== 4 BUDGET HEALTH
    s = d.slide(); d.brand(s); d.footer(s, 4)
    d.breadcrumb(s, "03 · Deep Dive: Budget Health")
    d.title(s, "Budget Health")
    d.subtitle(s, f"FY budget distribution by category · €{summ.get('budget_total_m','-')}M committed")
    budget = data.get("budget", [])
    bx, by, bw = 0.94, 4.2, 9.0
    maxpct = max((b.get("pct", 0) for b in budget), default=1) or 1
    for i, b in enumerate(budget):
        y = by + i * 1.15
        d.p0(d.tb(s, bx, y, 3.2, 0.4), b.get("legend_label", b.get("label", "")), size=13, bold=True, color=INK)
        barw = max(0.15, bw * b.get("pct", 0) / maxpct)
        d.rect(s, bx + 3.3, y + 0.02, barw, 0.5, _hex(b.get("color")))
        d.p0(d.tb(s, bx + 3.3 + barw + 0.2, y, 3.0, 0.5, anchor=MSO_ANCHOR.MIDDLE),
             f"€{b.get('value_m','-')}M  ·  {b.get('pct','-')}%", size=14, bold=True, color=INK)
    d.card_head(s, 13.3, 4.2, 5.76, 4.7, LCARD, accent=GREEN)
    tf = d.tb(s, 13.52, 4.4, 5.35, 0.4)
    d.p0(tf, "BUDGET INSIGHTS", size=13, bold=True, color=DGREEN)
    tf = d.tb(s, 13.52, 4.95, 5.35, 3.8)
    for i, ins in enumerate(data.get("budget_insights", [])):
        d.runs(tf, [(ins.get("bold", ""), True, INK), (ins.get("tail", ""), False, GRAY)],
               12.5, first=(i == 0), sa=10)

    # ===================================================== 5 LEADERSHIP DECISIONS
    s = d.slide(); d.brand(s); d.footer(s, 5)
    d.breadcrumb(s, "04 · Leadership Decisions")
    d.title(s, "Leadership Decisions")
    d.subtitle(s, data.get("actionItems", {}).get("subtitle", ""))
    cards = data.get("decision_cards", [])[:4]
    cw, ch = 9.0, 2.75
    for idx, c in enumerate(cards):
        col = idx % 2; row = idx // 2
        l = 0.94 + col * (cw + 0.12)
        t = 4.05 + row * (ch + 0.25)
        sev = SEV.get(c.get("severity", "").upper(), MUTE)
        d.card_head(s, l, t, cw, ch, LCARD, accent=sev)
        tf = d.tb(s, l + 0.25, t + 0.16, cw - 0.5, 0.35)
        d.runs(tf, [(f"{c.get('code','')}  ", True, sev),
                    (c.get("severity", ""), True, sev),
                    (f"   {c.get('title','')}", True, INK)], 13, first=True)
        tf = d.tb(s, l + 0.25, t + 0.62, cw - 0.5, 0.9)
        d.p0(tf, c.get("hero", ""), size=14, bold=True, color=BLACK)
        tf = d.tb(s, l + 0.25, t + 1.55, cw - 0.5, 0.7)
        d.p0(tf, c.get("context", ""), size=11.5, color=GRAY)
        tf = d.tb(s, l + 0.25, t + ch - 0.5, cw - 0.5, 0.35)
        d.runs(tf, [("Owner: ", True, MUTE), (c.get("owner", "-"), False, GRAY),
                    ("    Target: ", True, MUTE), (c.get("target", "-"), False, GRAY)], 11, first=True)

    # ===================================================== 6 PROJECTS AT RISK
    s = d.slide(); d.brand(s); d.footer(s, 6)
    d.breadcrumb(s, "05 · Projects at Risk")
    d.title(s, "Projects at Risk")
    d.subtitle(s, f"Top {len(data.get('at_risk_top5', []))} by RAG severity · €{summ.get('exposure_m','-')}M total exposure across {summ.get('rag_count','-')} flagged")
    _table(d, s,
           ["RAG", "Project", "Owner", "Target", "CY Budget", "%"],
           [[p.get("rag"), p.get("project"), p.get("owner"), p.get("go_live"),
             p.get("budget"), f"{p.get('pct','-')}%"] for p in data.get("at_risk_top5", [])],
           0.94, 4.1, 18.12, [1.3, 8.5, 2.6, 2.0, 2.3, 1.42], rag_col=0)

    # ===================================================== 7 KEY RISKS
    s = d.slide(); d.brand(s); d.footer(s, 7)
    d.breadcrumb(s, "06 · Key Risks")
    d.title(s, "Key Risks")
    d.subtitle(s, "Top portfolio risks, mitigations and owners")
    _table(d, s,
           ["Risk", "Severity", "Status", "Mitigation", "Owner", "Target"],
           [[r.get("description"), r.get("severity"), r.get("status"),
             r.get("mitigation"), r.get("owner"), r.get("target")] for r in data.get("key_risks", [])],
           0.94, 4.1, 18.12, [6.6, 1.7, 1.9, 4.6, 2.0, 1.32], sev_col=1, rowh=1.0)

    # ===================================================== 8 APPENDIX A.1
    s = d.slide(); d.brand(s); d.footer(s, 8)
    d.breadcrumb(s, "Appendix A.1 · Red / Amber projects")
    d.title(s, "Flagged Projects")
    projects = data.get("projects", [])
    d.subtitle(s, f"All {len(projects)} Red / Amber projects this period")
    _table(d, s,
           ["RAG", "Code", "Project", "Owner", "Go-live", "Budget", "%"],
           [[p.get("rag"), p.get("code"), p.get("project"), p.get("owner"),
             p.get("go_live"), p.get("budget"), f"{p.get('pct','-')}%"] for p in projects],
           0.94, 4.0, 18.12, [1.1, 1.7, 7.5, 2.5, 1.9, 2.0, 1.42], rag_col=0,
           rowh=max(0.42, min(0.55, 5.6 / max(len(projects), 1))), fsize=10)

    return d.prs


def _table(d, s, headers, rows, l, t, w, col_w, rag_col=None, sev_col=None, rowh=0.62, fsize=11.5):
    from pptx.util import Inches as I
    nr, nc = len(rows) + 1, len(headers)
    shp = s.shapes.add_table(nr, nc, I(l), I(t), I(w), I(rowh * nr))
    tbl = shp.table; tbl.first_row = True; tbl.horz_banding = False
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = I(cw)
    for i, h in enumerate(headers):
        c = tbl.cell(0, i); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = DKCARD
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = I(0.12); c.margin_top = I(0.04); c.margin_bottom = I(0.04)
        for p in c.text_frame.paragraphs:
            Deck._st(p, fsize, True, GREEN if i > 0 else WHITE, PP_ALIGN.LEFT)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci); c.text = str(val if val is not None else "-")
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LCARD
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = I(0.12); c.margin_top = I(0.03); c.margin_bottom = I(0.03)
            color, bold = GRAY, False
            if ci == rag_col and str(val) in RAG:
                color, bold = RAG[str(val)], True
            elif ci == sev_col and str(val).upper() in SEV:
                color, bold = SEV[str(val).upper()], True
            elif ci == 0:
                color, bold = INK, True
            for p in c.text_frame.paragraphs:
                Deck._st(p, fsize, bold, color, PP_ALIGN.LEFT)
    return tbl


def _data_file(ymd: str) -> Path | None:
    """Find a period's dashboard JSON in public/data/ or the csv backup folder."""
    for cand in (DATA / f"{ymd}.json", PUBLIC / "data_csv_backup" / f"{ymd}.json"):
        if cand.exists():
            return cand
    return None


def build(ymd: str) -> Path:
    """Render the deck for one period (ymd like '2026-01') and write it to public/."""
    jf = _data_file(ymd)
    if jf is None:
        raise FileNotFoundError(f"no dashboard data for {ymd}")
    data = json.loads(jf.read_text(encoding="utf-8"))
    prs = render(data)
    disp = data.get("period_display", ymd).replace(" ", "")
    out = PUBLIC / f"PortfolioIQ_{disp}.pptx"
    try:
        prs.save(str(out))
        return out
    except PermissionError:
        # Canonical file is locked (e.g. open in PowerPoint). Return a temp copy so the
        # download still works instead of failing the whole request.
        import tempfile
        tmp = Path(tempfile.gettempdir()) / f"PortfolioIQ_{disp}.pptx"
        prs.save(str(tmp))
        return tmp


if __name__ == "__main__":
    ymd = sys.argv[1] if len(sys.argv) > 1 else "2026-01"
    p = build(ymd)
    print(f"BUILT {p}  ({p.stat().st_size:,} bytes)")
